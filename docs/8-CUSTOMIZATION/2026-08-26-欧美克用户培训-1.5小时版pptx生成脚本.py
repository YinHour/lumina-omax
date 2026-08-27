#!/usr/bin/env python3
"""生成《把没用过的新能力一次演示给你看》1.5h 功能演示版 PPT（28 页 16:9）"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

INDIGO = RGBColor(0x43, 0x38, 0xCA)
INDIGO_DARK = RGBColor(0x31, 0x2E, 0x81)
INDIGO_DEEP = RGBColor(0x1E, 0x1B, 0x4B)
INDIGO_LIGHT = RGBColor(0xEE, 0xF2, 0xFF)
AMBER = RGBColor(0xD9, 0x77, 0x06)
AMBER_LIGHT = RGBColor(0xFD, 0xF3, 0xE7)
TEXT = RGBColor(0x1F, 0x29, 0x37)
SUB = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x05, 0x96, 0x69)
YAHEI = "Microsoft YaHei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def set_run(run, size=16, bold=False, color=TEXT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = YAHEI
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", YAHEI)


def add_box(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.text_frame.word_wrap = True
    return box


def para(tf, text, size=16, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
         space_after=6, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    set_run(run, size, bold, color)
    return p


def add_rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1.2)
    sp.shadow.inherit = False
    return sp


def rect_text(sp, text, size=13, bold=False, color=WHITE, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE, space_after=2):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    set_run(run, size, bold, color)
    return tf


def blank_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    return slide


def deep_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = INDIGO_DEEP
    return slide


def header(slide, title, subtitle=None, tag=None):
    box = add_box(slide, Inches(0.7), Inches(0.42), Inches(12.0), Inches(0.75))
    para(box.text_frame, title, 26, True, INDIGO_DARK, first=True)
    add_rect(slide, Inches(0.72), Inches(1.14), Inches(0.9), Pt(3.2), AMBER,
             shape=MSO_SHAPE.RECTANGLE)
    if subtitle:
        sbox = add_box(slide, Inches(0.7), Inches(1.28), Inches(11.0), Inches(0.4))
        para(sbox.text_frame, subtitle, 13.5, False, SUB, first=True)
    if tag:
        t = add_rect(slide, Inches(11.6), Inches(0.5), Inches(1.2), Inches(0.42), AMBER)
        rect_text(t, tag, 13, True, WHITE)


def footer(slide, text=None):
    fbox = add_box(slide, Inches(0.7), Inches(7.08), Inches(12.0), Inches(0.32))
    para(fbox.text_frame, text or "Lumiton·Omax 用户培训 · 1.5 小时功能演示版", 10,
         False, SUB, align=PP_ALIGN.RIGHT, first=True)


def section_slide(num, title, desc):
    s = deep_slide()
    nbox = add_box(s, Inches(0.9), Inches(2.15), Inches(2.0), Inches(1.2))
    para(nbox.text_frame, num, 52, True, AMBER, first=True)
    tbox = add_box(s, Inches(0.9), Inches(3.3), Inches(11.5), Inches(1.2))
    para(tbox.text_frame, title, 38, True, WHITE, first=True)
    dbox = add_box(s, Inches(0.9), Inches(4.55), Inches(11.0), Inches(1.0))
    para(dbox.text_frame, desc, 16, False, RGBColor(0xC7, 0xD2, 0xFE), first=True)
    return s


def bullets(slide, x, y, w, h, items, size=14.5, gap=8, title=None, tsize=16.5):
    box = add_box(slide, x, y, w, h)
    tf = box.text_frame
    first = True
    if title:
        para(tf, title, tsize, True, INDIGO_DARK, first=True, space_after=8)
        first = False
    for it in items:
        para(tf, "▪ " + it, size, False, TEXT, space_after=gap, first=first)
        first = False
    return box


# ============ P1 封面 ============
s = deep_slide()
kb = add_box(s, Inches(0.9), Inches(1.05), Inches(11.5), Inches(0.5))
para(kb.text_frame, "LUMITON·OMAX 用户培训", 18, True, RGBColor(0xA5, 0xB4, 0xFC), first=True)
tb = add_box(s, Inches(0.9), Inches(2.1), Inches(11.6), Inches(1.9))
para(tb.text_frame, "把「没用过的新能力」", 44, True, WHITE, first=True, space_after=4)
para(tb.text_frame, "一次演示给你看", 44, True, WHITE)
sub = add_box(s, Inches(0.9), Inches(4.1), Inches(11.6), Inches(0.6))
para(sub.text_frame, "—— 系统主要功能演示 · 聚焦 6 月以来新增的 30+ 项能力", 20, False,
     RGBColor(0xC7, 0xD2, 0xFE), first=True)
add_rect(s, Inches(0.92), Inches(4.9), Inches(1.4), Pt(3), AMBER, shape=MSO_SHAPE.RECTANGLE)
info = add_box(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.8))
para(info.text_frame, "对象：主管领导 + 系统使用人员（约 12 人）", 14, False,
     RGBColor(0xA5, 0xB4, 0xFC), first=True, space_after=2)
para(info.text_frame, "形式：系统实操演示为主 · 90 分钟 · 2026 年 8 月", 14, False,
     RGBColor(0xA5, 0xB4, 0xFC))

# ============ P2 议程 ============
s = blank_slide()
header(s, "90 分钟议程", "按使用顺序：登录 → 来源 → 聊天 → 研究 → 组织 → 安全")
agenda = [
    ("1", "开场：新功能地图", "5'", "全员", INDIGO),
    ("2", "账号与工作台：个人资料 · 只看我的", "5'", "全员", INDIGO),
    ("3", "来源管理新能力：筛选引用范围 · 下载 · 图片源", "13'", "使用人员", INDIGO),
    ("4", "聊天新体验：导览卡片 · 建议 · 等待状态", "13'", "使用人员", INDIGO),
    ("5", "科研 Agent：自主检索 + 10 个策研技能", "20'", "全员核心", AMBER),
    ("6", "全局「提问与搜索」：覆盖统计 · 历史", "10'", "使用人员", INDIGO),
    ("7", "笔记本组织：聚合笔记本 · 密码管理", "8'", "领导", INDIGO),
    ("8", "数据安全：出网脱敏透明网关", "5'", "领导", AMBER),
    ("9", "收尾：帮助中心 · 路线图 · 答疑", "11'", "全员", INDIGO),
]
y = 1.62
for tag, name, dur, who, color in agenda:
    add_rect(s, Inches(0.7), Inches(y), Inches(0.55), Inches(0.52), color)
    rect_text(s.shapes[-1], tag, 15, True, WHITE)
    nb = add_box(s, Inches(1.45), Inches(y + 0.02), Inches(9.2), Inches(0.5))
    para(nb.text_frame, name, 14.5, True, TEXT, first=True)
    db = add_box(s, Inches(10.7), Inches(y + 0.04), Inches(2.0), Inches(0.45))
    para(db.text_frame, dur + " · " + who, 11.5, False, SUB, first=True)
    y += 0.6
footer(s)

# ============ P3 新功能地图 ============
s = blank_slide()
header(s, "新功能地图（本场覆盖）", None, tag="总览")
layers = [
    ("账号层", "个人资料页 · 只看我的 · 最近笔记本", False),
    ("来源层", "搜索分页 · 来源筛选(本次引用范围) · 仅已嵌入可选 · 数量上限 · 下载(MD/ZIP) · 独立图片源 · 删源三规则", True),
    ("聊天层", "导览卡片 · 回答后 3 条建议 · 阶段等待状态 · 超时气泡 · 保存笔记 · 联网搜索降级", True),
    ("研究层", "⭐科研 Agent（自主检索 + 10 个策研技能）· ⭐全局提问与搜索（覆盖统计 + 历史）", True),
    ("组织层", "聚合笔记本 · 密码全生命周期 · 来源数量上限", False),
    ("安全层", "⭐出网脱敏透明网关（库内原文不动 · 出网替换 · 回网还原）", True),
]
y = 1.75
for name, desc, highlight in layers:
    fill = AMBER if "⭐" in name or name == "安全层" else INDIGO
    hd = add_rect(s, Inches(0.7), Inches(y), Inches(1.35), Inches(0.68), fill)
    rect_text(hd, name, 13.5, True, WHITE)
    db = add_box(s, Inches(2.25), Inches(y + 0.05), Inches(10.4), Inches(0.62))
    para(db.text_frame, desc, 13.5, False, TEXT, first=True)
    y += 0.78
lg = add_rect(s, Inches(0.7), Inches(6.45), Inches(12.05), Inches(0.5), INDIGO_LIGHT, INDIGO)
rect_text(lg, "⭐ = 8 月新增（最新）　🆕 = 6-7 月新增（多数人没用过）　其余 = 已有能力补强", 13.5,
          True, INDIGO_DARK, PP_ALIGN.LEFT)
footer(s)

# ============ P4 章节 账号 ============
section_slide("01", "账号与工作台", "5 分钟 · 先说清楚自己的身份、数据和入口在哪里。")

# ============ P5 账号与工作台 ============
s = blank_slide()
header(s, "账号与工作台", None, tag="5'")
items_l = [
    "🆕 个人资料页 /profile：改显示名称、改密码（旧密码验证），无需找管理员",
    "🆕 只看我的：一键只看自己创建的笔记本，可叠加名称搜索",
    "🆕 最近笔记本：侧边栏快进入口，长标题 hover 查看全名",
    "词元统计：聊天输入框实时显示（Token 已统一为「词元」）",
]
bullets(s, 0.7, 1.9, 12.0, 4.2, items_l, size=16, gap=14)
bb = add_rect(s, Inches(0.7), Inches(6.35), Inches(12.05), Inches(0.6), INDIGO_LIGHT, INDIGO)
rect_text(bb, "现场操作：打开个人资料页 → 笔记本页演示「只看我的」", 14, True, INDIGO_DARK,
          PP_ALIGN.LEFT)
footer(s)

# ============ P6 章节 来源 ============
section_slide("02", "来源管理新能力", "13 分钟 · 资料多了以后：找得到、选得准、交得出。")

# ============ P7 来源：搜索分页 + 筛选 ============
s = blank_slide()
header(s, "来源管理 ①：搜索分页 + 筛选引用范围", None, tag="13'")
lbox = add_box(s, Inches(0.7), Inches(1.75), Inches(5.9), Inches(4.6))
para(lbox.text_frame, "搜索与分页（Sources 全局页）", 16.5, True, INDIGO_DARK, first=True,
     space_after=8)
for it in ["标题搜索（回车触发）+ 总条数/页码",
           "搜索无结果有明确提示，不再像系统没数据",
           "上传人列：多人共用时追溯资料来源",
           "添加现有来源弹窗支持搜索 + 全选筛选结果"]:
    para(lbox.text_frame, "▪ " + it, 14, False, TEXT, space_after=8)
rbox = add_box(s, Inches(6.85), Inches(1.75), Inches(5.9), Inches(4.6))
para(rbox.text_frame, "来源筛选 = 本次聊天引用范围（核心）", 16.5, True, INDIGO_DARK,
     first=True, space_after=8)
for it in ["打开笔记本 → 来源栏筛选框",
           "筛选结果支持「全选 / 取消全选」",
           "资料属于项目 ≠ 本次回答引用哪些",
           "取消勾选只是不参与本次回答，来源仍在笔记本"]:
    para(rbox.text_frame, "▪ " + it, 14, False, TEXT, space_after=8)
bb = add_rect(s, Inches(0.7), Inches(6.35), Inches(12.05), Inches(0.6), AMBER_LIGHT, AMBER)
rect_text(bb, "现场演示：筛选「HX」→ 全选 → 只让筛选结果参与本次聊天（30 秒）", 14, True,
          RGBColor(0x92, 0x40, 0x0E), PP_ALIGN.LEFT)
footer(s)

# ============ P8 来源：下载 + 图片源 ============
s = blank_slide()
header(s, "来源管理 ②：解析结果可交付 + 图片也能查", None, tag="13'")
lbox = add_box(s, Inches(0.7), Inches(1.75), Inches(5.9), Inches(4.6))
para(lbox.text_frame, "解析结果下载（来源详情页）", 16.5, True, INDIGO_DARK, first=True,
     space_after=8)
for it in ["Markdown 下载：UTF-8 全文，可直接编辑归档",
           "ZIP 包下载：Markdown + 全部抽取图片",
           "图片路径已改写为相对路径，解压即可预览带图文档",
           "重要解析结果可以交付、复查、分享"]:
    para(lbox.text_frame, "▪ " + it, 14, False, TEXT, space_after=8)
rbox = add_box(s, Inches(6.85), Inches(1.75), Inches(5.9), Inches(4.6))
para(rbox.text_frame, "独立图片源 + Vision 描述", 16.5, True, INDIGO_DARK, first=True,
     space_after=8)
for it in ["上传图片即成为来源：先显示原图，再显示 AI 描述",
           "TIFF 无法在浏览器预览 → 自动转 PNG",
           "图片描述纳入向量搜索：按图片内容语义搜索",
           "文档内嵌图（PDF/Excel）也自动抽取并描述"]:
    para(rbox.text_frame, "▪ " + it, 14, False, TEXT, space_after=8)
bb = add_rect(s, Inches(0.7), Inches(6.35), Inches(12.05), Inches(0.6), INDIGO_LIGHT, INDIGO)
rect_text(bb, "现场演示：下载一个 ZIP 包 → 打开一个图片源（1 分钟）", 14, True, INDIGO_DARK,
          PP_ALIGN.LEFT)
footer(s)

# ============ P9 来源：三规则 ============
s = blank_slide()
header(s, "来源管理 ③：三条防误操作规则", None, tag="13'")
cards = [
    ("仅已嵌入可选", "添加现有来源弹窗只列出处理完成的来源；未处理完的不出现、不显示禁选原因，不误选「半成品」"),
    ("来源数量上限", "单笔记本来源数默认 50，管理员在「设置 → 文件管理」可调 1-200；达到上限有明确提示，防止资料无限堆积"),
    ("删源三规则", "自己的 + 未被引用 → 直接删；别人的 → 只能「移除」；自己的 + 多笔记本引用 → 需管理员密码（后端校验）"),
]
x = 0.7
for title, body in cards:
    blk = add_rect(s, Inches(x), Inches(1.9), Inches(3.95), Inches(3.6), INDIGO_LIGHT, INDIGO)
    tf = blk.text_frame
    tf.margin_left = Pt(14)
    tf.margin_top = Pt(14)
    para(tf, title, 16.5, True, INDIGO_DARK, first=True, space_after=10)
    para(tf, body, 13, False, TEXT)
    x += 4.05
footer(s)

# ============ P10 章节 聊天 ============
section_slide("03", "聊天新体验", "13 分钟 · 从「不会问」到「被引导着问」。")

# ============ P11 聊天：导览 + 建议 ============
s = blank_slide()
header(s, "聊天 ①：导览卡片 + 回答后的下一步建议", None, tag="13'")
c1 = add_rect(s, Inches(0.7), Inches(1.85), Inches(5.95), Inches(4.0), INDIGO_LIGHT, INDIGO)
tf = c1.text_frame
tf.margin_left = Pt(16)
tf.margin_top = Pt(16)
para(tf, "🆕 导览卡片（首次导入后自动生成）", 17, True, INDIGO_DARK, first=True, space_after=10)
for it in ["笔记本摘要 + 关键要点",
           "3 条建议问题，点击即直接发送",
           "不知道问什么？从导览卡片开始"]:
    para(tf, "▪ " + it, 14, False, TEXT, space_after=8)
c2 = add_rect(s, Inches(6.8), Inches(1.85), Inches(5.95), Inches(4.0), INDIGO_LIGHT, INDIGO)
tf = c2.text_frame
tf.margin_left = Pt(16)
tf.margin_top = Pt(16)
para(tf, "🆕 每轮回答后的 3 条下一步建议", 17, True, INDIGO_DARK, first=True, space_after=10)
for it in ["结合本轮问题 + 主回答 + 笔记本上下文",
           "主回答完成即恢复输入框，建议异步生成不阻塞",
           "点击建议直接追问，连续对话不跑偏"]:
    para(tf, "▪ " + it, 14, False, TEXT, space_after=8)
bb = add_rect(s, Inches(0.7), Inches(6.15), Inches(12.05), Inches(0.6), AMBER_LIGHT, AMBER)
rect_text(bb, "两条内置引导 = 把「许愿式提问」拉回研发闭环", 14.5, True,
          RGBColor(0x92, 0x40, 0x0E), PP_ALIGN.LEFT)
footer(s)

# ============ P12 聊天：等待状态 + 超时 ============
s = blank_slide()
header(s, "聊天 ②：慢的时候，你能看到系统在干什么", None, tag="13'")
steps = ["获取上下文", "联网搜索", "等待模型响应（N 秒）", "模型已开始输出", "正在生成建议"]
x = 0.7
w = 2.35
for i, st in enumerate(steps):
    fill = INDIGO if i < 4 else AMBER
    blk = add_rect(s, Inches(x), Inches(1.85), Inches(w), Inches(0.75), fill)
    rect_text(blk, f"{i+1}. {st}", 13.5, True, WHITE)
    if i < 4:
        ab = add_box(s, Inches(x + w), Inches(2.03), Inches(0.15), Inches(0.35))
        para(ab.text_frame, "›", 16, True, SUB, PP_ALIGN.CENTER, first=True)
    x += w + 0.15
lbox = add_box(s, Inches(0.7), Inches(3.0), Inches(12.0), Inches(3.2))
para(lbox.text_frame, "超时保护（不再像卡死）", 16.5, True, INDIGO_DARK, first=True,
     space_after=8)
for it in ["等待秒数实时心跳：区分「工作中」还是「卡死」",
           "超时以 AI 气泡提示（中文操作指引）：缩小来源范围 / 新建会话重试，不再弹英文 toast",
           "移动端切 Tab 不中断回答；停止生成保留已生成内容"]:
    para(lbox.text_frame, "▪ " + it, 14.5, False, TEXT, space_after=8)
footer(s)

# ============ P13 聊天：保存笔记 + 联网 ============
s = blank_slide()
header(s, "聊天 ③：保存为笔记 + 联网搜索", None, tag="13'")
lbox = add_box(s, Inches(0.7), Inches(1.85), Inches(5.9), Inches(4.2))
para(lbox.text_frame, "保存为笔记（回答沉淀）", 16.5, True, INDIGO_DARK, first=True,
     space_after=8)
for it in ["每轮有价值回答一键存为笔记",
           "演示笔记本已存 7 份递进式笔记（问题定义卡→管理摘要）",
           "笔记参与搜索与后续聊天上下文",
           "Ask 答案同样可复制 / 保存为笔记"]:
    para(lbox.text_frame, "▪ " + it, 14, False, TEXT, space_after=8)
rbox = add_box(s, Inches(6.85), Inches(1.85), Inches(5.9), Inches(4.2))
para(rbox.text_frame, "联网搜索（Tavily）", 16.5, True, INDIGO_DARK, first=True, space_after=8)
for it in ["Chat 面板开关，AI 可检索实时网页",
           "查询有超时控制 + 失败降级，不会无限等待",
           "本地来源与网络引用分开编号",
           "「搜索次数用完」≠ 模型配额，是 Tavily 上限"]:
    para(rbox.text_frame, "▪ " + it, 14, False, TEXT, space_after=8)
footer(s)

# ============ P14 章节 科研 Agent ============
section_slide("04", "科研 Agent", "20 分钟 · 全场核心：把复杂问题交给它，它自己检索、阅读、综合。")

# ============ P15 什么是科研 Agent ============
s = blank_slide()
header(s, "什么是科研 Agent：聊天页切「科研」模式", None, tag="核心")
rows = [
    ("", "快速聊天（Quick）", "科研 Agent（Research）"),
    ("上下文", "你手动选来源、定模式", "它自主检索：列源 → 向量检索 → 分段读取 → 再检索"),
    ("提问方式", "逐步递进的标准提问", "只问一个复杂问题，它自主规划检索路径"),
    ("能力", "基于选定上下文回答", "工具调用 + 联网 + 科学数据库 + 10 个策研技能"),
    ("适用", "已知来源、精读分析", "证据脉络不明、跨文档诊断、复杂综合"),
]
y = 1.8
for label, a, b in rows:
    if label:
        lb = add_rect(s, Inches(0.7), Inches(y), Inches(1.5), Inches(0.62), INDIGO_LIGHT,
                      INDIGO)
        rect_text(lb, label, 13.5, True, INDIGO_DARK)
    else:
        lb = add_rect(s, Inches(0.7), Inches(y), Inches(1.5), Inches(0.62), INDIGO_DEEP)
        rect_text(lb, "", 13.5, True, WHITE)
    ca = add_box(s, Inches(2.35), Inches(y + 0.04), Inches(5.2), Inches(0.6))
    para(ca.text_frame, a, 13, False, SUB, first=True)
    cb = add_rect(s, Inches(7.75), Inches(y), Inches(5.0), Inches(0.62), AMBER_LIGHT, AMBER)
    rect_text(cb, b, 12.5, True, RGBColor(0x92, 0x40, 0x0E), PP_ALIGN.LEFT)
    y += 0.78
footer(s)

# ============ P16 能力清单 ============
s = blank_slide()
header(s, "科研 Agent 的能力清单", None, tag="核心")
lbox = add_box(s, Inches(0.7), Inches(1.8), Inches(5.9), Inches(4.4))
para(lbox.text_frame, "内置工具（只读，全程限笔记本范围）", 16.5, True, INDIGO_DARK,
     first=True, space_after=8)
for it in ["列出来源与笔记清单",
           "语义检索笔记本证据",
           "按 ID 分段读取来源 / 笔记",
           "跨笔记本发现（需显式授权）",
           "联网搜索 · 科学数据库连接器"]:
    para(lbox.text_frame, "▪ " + it, 14, False, TEXT, space_after=8)
rbox = add_box(s, Inches(6.85), Inches(1.8), Inches(5.9), Inches(4.4))
para(rbox.text_frame, "10 个策研技能（按问题自动加载）", 16.5, True, INDIGO_DARK, first=True,
     space_after=8)
skills = ["化学身份与物性 · 竞争假设 · DOE 统计计划",
          "证据批判性评估 · 配方兼容矩阵",
          "高温高压盐水验证 · 文献 DOI 验证",
          "油井水泥外加剂诊断 · 放大验证 Gate",
          "结构化研究报告"]
for it in skills:
    para(rbox.text_frame, "▪ " + it, 13, False, TEXT, space_after=6)
footer(s)

# ============ P17 真实工作轨迹 ============
s = blank_slide()
header(s, "真实案例：一个问题的自主工作轨迹", "素材：演示笔记本「科研Agent示范_大温差缓凝剂顶部强度」（10 份源）", tag="核心")
traj = [
    ("1", "第 1 轮", "列出来源清单，决定先读时间线早期文档"),
    ("2", "第 2-3 轮", "逐份读取研发报告（每份分段，共 10 份）"),
    ("3", "第 4-5 轮", "检索「顶部强度/缓凝」矛盾证据，再读关键段落"),
    ("4", "第 6 轮", "停止检索，进入最终综合"),
    ("5", "自主纠错", "发现用了猜测的来源 ID → 改用工具返回的真实 ID 列表"),
]
y = 1.95
for tag, name, desc in traj:
    fill = AMBER if "纠错" in name else INDIGO
    add_rect(s, Inches(0.7), Inches(y), Inches(0.7), Inches(0.62), fill)
    rect_text(s.shapes[-1], tag, 14, True, WHITE)
    nb = add_box(s, Inches(1.6), Inches(y + 0.02), Inches(3.0), Inches(0.58))
    para(nb.text_frame, name, 14.5, True, TEXT, first=True)
    db = add_box(s, Inches(4.7), Inches(y + 0.02), Inches(8.0), Inches(0.58))
    para(db.text_frame, desc, 13.5, False, TEXT, first=True)
    y += 0.8
bb = add_rect(s, Inches(0.7), Inches(6.15), Inches(12.05), Inches(0.6), AMBER_LIGHT, AMBER)
rect_text(bb, "全程 6 轮迭代、10+ 次工具调用、约 3 分钟 —— 过程可观测、可审计（日志留痕）", 14,
          True, RGBColor(0x92, 0x40, 0x0E), PP_ALIGN.LEFT)
footer(s)

# ============ P18 成果 1 ============
s = blank_slide()
header(s, "成果 ①：证据脉络综合报告（170 秒自动产出）", None, tag="核心")
c1 = add_rect(s, Inches(0.7), Inches(1.85), Inches(5.95), Inches(3.9), INDIGO_LIGHT, INDIGO)
tf = c1.text_frame
tf.margin_left = Pt(16)
tf.margin_top = Pt(16)
para(tf, "完整时间线表", 16.5, True, INDIGO_DARK, first=True, space_after=8)
for it in ["2025.6 → 2026.8 共 9 个时间节点",
           "每节点：关键发现 + 遗留问题 + 来源 ID",
           "四条研发路线：新单体 → QS-20S → 咪唑/1.6SG → 大温差 HX-36LB"]:
    para(tf, "▪ " + it, 13.5, False, TEXT, space_after=8)
c2 = add_rect(s, Inches(6.8), Inches(1.85), Inches(5.95), Inches(3.9), AMBER_LIGHT, AMBER)
tf = c2.text_frame
tf.margin_left = Pt(16)
tf.margin_top = Pt(16)
para(tf, "四类矛盾证据（诊断价值）", 16.5, True, RGBColor(0x92, 0x40, 0x0E), first=True,
     space_after=8)
for it in ["强缓凝 → 顶部强度滞后（90℃ 4 天未初凝）",
           "温度倒挂：升温 5℃ 稠化时间反而翻倍",
           "降黏与缓凝时长的跷跷板",
           "渤星标杆反衬（QS-20S 压低 40℃ 强度）"]:
    para(tf, "▪ " + it, 13.5, False, TEXT, space_after=8)
bb = add_rect(s, Inches(0.7), Inches(6.15), Inches(12.05), Inches(0.6), INDIGO_LIGHT, INDIGO)
rect_text(bb, "严格区分「证据」与「推断」：原始实测 [source] vs 转述派生 [source_insight] 标注待核对", 14,
          True, INDIGO_DARK, PP_ALIGN.LEFT)
footer(s)

# ============ P19 成果 2 ============
s = blank_slide()
header(s, "成果 ②：机理假设 + 6 组验证方案（248 秒自动产出）", None, tag="核心")
lbox = add_box(s, Inches(0.7), Inches(1.85), Inches(5.9), Inches(4.2))
para(lbox.text_frame, "三个机理假设", 16.5, True, INDIGO_DARK, first=True, space_after=8)
for it in ["H1（优先）：缓凝剂在 90℃ 低温段过度抑制顶部水化诱导期",
           "H2：降温路径诱发吸附迟滞/构象变化（与倒挂同源）",
           "H3：降黏/分散组分的低温竞争吸附"]:
    para(lbox.text_frame, "▪ " + it, 13.5, False, TEXT, space_after=8)
rbox = add_box(s, Inches(6.85), Inches(1.85), Inches(5.9), Inches(4.2))
para(rbox.text_frame, "H1 六组验证方案", 16.5, True, INDIGO_DARK, first=True, space_after=8)
for it in ["空白基线 → 渤星标杆 → 缺陷复现",
           "→ 剂量减半 → 成核封锁探针 → 机制定量",
           "每组含通过标准 + 失败判据",
           "附「组 4/5 失败 → 转 H2/H3」的判定逻辑"]:
    para(rbox.text_frame, "▪ " + it, 13.5, False, TEXT, space_after=8)
bb = add_rect(s, Inches(0.7), Inches(6.35), Inches(12.05), Inches(0.6), AMBER_LIGHT, AMBER)
rect_text(bb, "元发现：现有数据几乎全部是单次实验 —— 方案要求每组 ≥2 次平行", 14, True,
          RGBColor(0x92, 0x40, 0x0E), PP_ALIGN.LEFT)
footer(s)

# ============ P20 适用场景 ============
s = blank_slide()
header(s, "什么时候用科研 Agent", None, tag="核心")
lbox = add_box(s, Inches(0.7), Inches(1.9), Inches(5.9), Inches(4.0))
para(lbox.text_frame, "用科研 Agent（研究模式）", 17, True, INDIGO_DARK, first=True,
     space_after=10)
for it in ["证据脉络不明，需要先盘点",
           "跨多份文档诊断、比较、综合",
           "复杂问题一次交付结构化报告",
           "「你只问问题，它负责检索与阅读」"]:
    para(lbox.text_frame, "▪ " + it, 14.5, False, TEXT, space_after=9)
rbox = add_box(s, Inches(6.85), Inches(1.9), Inches(5.9), Inches(4.0))
para(rbox.text_frame, "用快速聊天（快速模式）", 17, True, INDIGO_DARK, first=True,
     space_after=10)
for it in ["已知要参考哪些来源",
           "精读分析、逐步追问验证",
           "上下文敏感的单点问题",
           "「你控制证据，AI 负责回答」"]:
    para(rbox.text_frame, "▪ " + it, 14.5, False, TEXT, space_after=9)
bb = add_rect(s, Inches(0.7), Inches(6.35), Inches(12.05), Inches(0.6), AMBER_LIGHT, AMBER)
rect_text(bb, "现场演示：研究模式发一个短问题（约 2-4 分钟，模型不稳时展示成品报告）", 14, True,
          RGBColor(0x92, 0x40, 0x0E), PP_ALIGN.LEFT)
footer(s)

# ============ P21 章节 全局提问 ============
section_slide("05", "全局「提问与搜索」", "10 分钟 · 这个页面 8 月才修好，多数人可能从没用成功过。")

# ============ P22 全局提问 ============
s = blank_slide()
header(s, "全局提问与搜索：覆盖统计 + 历史记录", None, tag="10'")
lbox = add_box(s, Inches(0.7), Inches(1.85), Inches(5.9), Inches(4.3))
para(lbox.text_frame, "检索覆盖统计（每个回答附带）", 16.5, True, INDIGO_DARK, first=True,
     space_after=8)
for it in ["来源总数 / 可检索来源 / 本次命中来源",
           "一眼判断回答覆盖了多少资料",
           "命中少 = 该补资料或换问法",
           "最终回答不会把「命中数」误说成「总数」"]:
    para(lbox.text_frame, "▪ " + it, 14, False, TEXT, space_after=8)
rbox = add_box(s, Inches(6.85), Inches(1.85), Inches(5.9), Inches(4.3))
para(rbox.text_frame, "体验细节", 16.5, True, INDIGO_DARK, first=True, space_after=8)
for it in ["历史记录：浏览器保存问答+覆盖统计，刷新不丢",
           "思考过程与最终答案分开渲染",
           "可中途停止，保留已生成内容",
           "菜单与标题已统一为「提问与搜索」"]:
    para(rbox.text_frame, "▪ " + it, 14, False, TEXT, space_after=8)
bb = add_rect(s, Inches(0.7), Inches(6.35), Inches(12.05), Inches(0.6), AMBER_LIGHT, AMBER)
rect_text(bb, "8 月修复：此前「模型供应商返回错误」实为内部并行检索合并问题——与供应商无关，已修复", 14,
          True, RGBColor(0x92, 0x40, 0x0E), PP_ALIGN.LEFT)
footer(s)

# ============ P23 章节 笔记本组织 ============
section_slide("06", "笔记本组织", "8 分钟 · 领导关注：项目怎么分层、资料怎么收口。")

# ============ P24 聚合 + 密码 ============
s = blank_slide()
header(s, "聚合笔记本 + 密码全生命周期", None, tag="8'")
top = add_rect(s, Inches(0.7), Inches(1.8), Inches(12.05), Inches(0.95), INDIGO_DEEP)
tft = top.text_frame
tft.margin_left = Pt(14)
para(tft, "聚合笔记本：多子笔记本合并成一个视图 —— 不拷贝数据、毫秒级同步", 15.5, True, WHITE,
     first=True, space_after=3)
para(tft, "试点体系 = 1 聚合 + 6 子（产品说明 / 研发报告 / 实验记录 / 现场反馈 / 水泥批次 / 外部文献）",
     13.5, False, RGBColor(0xE0, 0xE7, 0xFF))
lbox = add_box(s, Inches(0.7), Inches(3.1), Inches(5.9), Inches(3.2))
para(lbox.text_frame, "密码全生命周期", 16.5, True, INDIGO_DARK, first=True, space_after=8)
for it in ["设密码 / 改密码 / 撤销密码，随时可调",
           "仅笔记本创建者可管理",
           "管理员密码可绕过（兜底）",
           "敏感项目、客户资料建议加密"]:
    para(lbox.text_frame, "▪ " + it, 14, False, TEXT, space_after=8)
rbox = add_box(s, Inches(6.85), Inches(3.1), Inches(5.9), Inches(3.2))
para(rbox.text_frame, "权限收敛", 16.5, True, INDIGO_DARK, first=True, space_after=8)
for it in ["三点菜单（密码/归档/删除）仅创建者可见",
           "来源上限默认 50，防止堆积",
           "只看我的 + 搜索 + 活动/聚合/归档分组"]:
    para(rbox.text_frame, "▪ " + it, 14, False, TEXT, space_after=8)
footer(s)

# ============ P25 章节 数据安全 ============
section_slide("07", "数据安全：出网脱敏", "5 分钟 · 领导关注：核心 IP 信息不出企业边界。")

# ============ P26 脱敏 ============
s = blank_slide()
header(s, "出网脱敏透明网关：使用无感", None, tag="5'")
top = add_rect(s, Inches(0.7), Inches(1.75), Inches(12.05), Inches(1.05), INDIGO_DEEP)
tft = top.text_frame
tft.margin_left = Pt(14)
para(tft, "库内原文不动 → 发往外部 AI 瞬间临时替换 → 返回自动还原", 16, True, WHITE,
     first=True, space_after=4)
para(tft, "公司名→某企业 · 井号→实验井A · 产品代号→减阻剂A · 电话→888888；您看到的界面与原来一致，检索、引用、存储全部用原文",
     13.5, False, RGBColor(0xE0, 0xE7, 0xFF))
notes = [
    ("注意 1：电话显示 888888", "规则固化不可还原；真实电话请查阅原始资料"),
    ("注意 2：联网搜索用科学主题词", "身份词以代号出网搜不到；用「高温下降失水剂性能」而非「张三的实验」"),
    ("注意 3：新人名需录词典", "井号/电话/产品代号自动识别；人名需管理员在「设置→内容脱敏」录入"),
]
x = 0.7
for title, body in notes:
    blk = add_rect(s, Inches(x), Inches(3.15), Inches(3.95), Inches(2.0), AMBER_LIGHT, AMBER)
    tf = blk.text_frame
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(12)
    para(tf, title, 14, True, RGBColor(0x92, 0x40, 0x0E), first=True, space_after=8)
    para(tf, body, 12.5, False, TEXT)
    x += 4.05
bb = add_rect(s, Inches(0.7), Inches(5.5), Inches(12.05), Inches(0.75), INDIGO_LIGHT, INDIGO)
bbtf = bb.text_frame
bbtf.margin_left = Pt(14)
para(bbtf, "边界如实说明：图片、音频不经过文本脱敏；嵌入向量与库内存储按信任边界决策不脱敏", 13.5,
     False, INDIGO_DARK, first=True)
footer(s)

# ============ P27 收尾 ============
s = blank_slide()
header(s, "收尾：帮助中心 · 路线图 · 练习素材", None, tag="11'")
lbox = add_box(s, Inches(0.7), Inches(1.8), Inches(5.9), Inches(4.5))
para(lbox.text_frame, "会后可用", 16.5, True, INDIGO_DARK, first=True, space_after=8)
for it in ["/help 帮助中心：全中文操作手册",
           "8 月已同步全部新功能说明",
           "两个演示笔记本留作练习：\n  ① HX-16L 证据闭环（快速聊天 7 阶段）\n  ② 大温差缓凝剂（科研 Agent）"]:
    para(lbox.text_frame, "▪ " + it, 14, False, TEXT, space_after=8)
rbox = add_box(s, Inches(6.85), Inches(1.8), Inches(5.9), Inches(4.5))
para(rbox.text_frame, "路线图（规划中，如实说明）", 16.5, True, RGBColor(0x92, 0x40, 0x0E),
     first=True, space_after=8)
for it in ["模型故障自动切换",
           "实验数据索引（Excel 逐条结构化）",
           "实验时间线（按产品追溯）",
           "产品图谱与版本迭代图",
           "数据备份恢复流程文档化"]:
    para(rbox.text_frame, "▪ " + it, 14, False, TEXT, space_after=8)
bb = add_rect(s, Inches(0.7), Inches(6.45), Inches(12.05), Inches(0.55), INDIGO_LIGHT, INDIGO)
rect_text(bb, "答疑：现场提问 + 会后群内跟进", 14.5, True, INDIGO_DARK, PP_ALIGN.LEFT)
footer(s)

# ============ P28 结尾 ============
s = deep_slide()
tb = add_box(s, Inches(0.9), Inches(2.5), Inches(11.6), Inches(1.6))
para(tb.text_frame, "把没用过的新能力，", 38, True, WHITE, first=True, space_after=6)
para(tb.text_frame, "变成每天在用的生产力。", 38, True, WHITE)
add_rect(s, Inches(0.92), Inches(4.4), Inches(1.4), Pt(3), AMBER, shape=MSO_SHAPE.RECTANGLE)
sb = add_box(s, Inches(0.9), Inches(4.8), Inches(11.6), Inches(1.0))
para(sb.text_frame, "今天之后：登录 → 传资料 → 筛选引用 → 聊天被引导 → 复杂问题交给科研 Agent", 18,
     False, RGBColor(0xC7, 0xD2, 0xFE), first=True)
qb = add_box(s, Inches(0.9), Inches(6.0), Inches(11.6), Inches(0.6))
para(qb.text_frame, "会后材料：/help 帮助中心 · 两个演示笔记本 · 本演示大纲", 14, False,
     RGBColor(0xA5, 0xB4, 0xFC), first=True)

OUT = "/Users/omax/YinShiApp/lumina-omax/docs/8-CUSTOMIZATION/2026-08-26-欧美克用户培训-1.5小时功能演示版.pptx"
prs.save(OUT)
print(f"saved: {OUT} ({len(prs.slides._sldIdLst)} slides)")
