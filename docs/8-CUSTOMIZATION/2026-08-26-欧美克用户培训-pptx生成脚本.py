#!/usr/bin/env python3
"""生成《让每一轮研发都可追踪、可复盘、可预测》用户培训 PPT（34 页 16:9）"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

INDIGO = RGBColor(0x43, 0x38, 0xCA)
INDIGO_DARK = RGBColor(0x31, 0x2E, 0x81)
INDIGO_DEEP = RGBColor(0x1E, 0x1B, 0x4B)
INDIGO_LIGHT = RGBColor(0xEE, 0xF2, 0xFF)
AMBER = RGBColor(0xD9, 0x77, 0x06)
AMBER_LIGHT = RGBColor(0xFD, 0xF3, 0xE7)
TEXT = RGBColor(0x1F, 0x29, 0x37)
SUB = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
YAHEI = "Microsoft YaHei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def set_run(run, size=16, bold=False, color=TEXT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = YAHEI
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", YAHEI)


def add_box(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.text_frame.word_wrap = True
    return box


def para(tf, text, size=16, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
         space_after=6, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    set_run(run, size, bold, color)
    return p


def add_rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1.2)
    sp.shadow.inherit = False
    return sp


def rect_text(sp, text, size=13, bold=False, color=WHITE, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE, space_after=2):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    set_run(run, size, bold, color)
    return tf


def blank_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    return slide


def deep_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = INDIGO_DEEP
    return slide


def header(slide, title, subtitle=None):
    add_box(slide, Inches(0.7), Inches(0.42), Inches(12.0), Inches(0.75)).text_frame
    box = add_box(slide, Inches(0.7), Inches(0.42), Inches(12.0), Inches(0.75))
    para(box.text_frame, title, 28, True, INDIGO_DARK, first=True)
    add_rect(slide, Inches(0.72), Inches(1.14), Inches(0.9), Pt(3.2), AMBER,
             shape=MSO_SHAPE.RECTANGLE)
    if subtitle:
        sbox = add_box(slide, Inches(0.7), Inches(1.28), Inches(12.0), Inches(0.4))
        para(sbox.text_frame, subtitle, 13.5, False, SUB, first=True)


def footer(slide, text=None):
    fbox = add_box(slide, Inches(0.7), Inches(7.08), Inches(12.0), Inches(0.32))
    para(fbox.text_frame, text or "Lumiton·Omax 用户培训 · 证据驱动的研发闭环", 10,
         False, SUB, align=PP_ALIGN.RIGHT, first=True)


def section_slide(num, title, desc):
    s = deep_slide()
    nbox = add_box(s, Inches(0.9), Inches(2.2), Inches(2.0), Inches(1.2))
    para(nbox.text_frame, num, 54, True, AMBER, first=True)
    tbox = add_box(s, Inches(0.9), Inches(3.35), Inches(11.5), Inches(1.2))
    para(tbox.text_frame, title, 40, True, WHITE, first=True)
    dbox = add_box(s, Inches(0.9), Inches(4.6), Inches(11.0), Inches(1.0))
    para(dbox.text_frame, desc, 16, False, RGBColor(0xC7, 0xD2, 0xFE), first=True)
    return s


def practice_slide(num, stage, steps, prompt_lines, output):
    s = blank_slide()
    header(s, f"阶段 {stage}")
    tag = add_rect(s, Inches(0.7), Inches(1.42), Inches(2.0), Inches(0.46), AMBER)
    rect_text(tag, f"现场实操 {num}", 14, True, WHITE)
    sbox = add_box(s, Inches(0.7), Inches(2.05), Inches(12.0), Inches(1.25))
    para(sbox.text_frame, "操作步骤", 14.5, True, INDIGO_DARK, first=True, space_after=4)
    marks = ["①", "②", "③", "④"]
    for i, st in enumerate(steps):
        para(sbox.text_frame, f"{marks[i]} {st}", 13.5, False, TEXT, space_after=3)
    qbox = add_rect(s, Inches(0.7), Inches(3.42), Inches(12.0), Inches(3.28),
                    INDIGO_LIGHT, INDIGO)
    qtf = qbox.text_frame
    qtf.margin_left = Pt(14)
    qtf.margin_right = Pt(14)
    qtf.margin_top = Pt(10)
    para(qtf, "标准提问（直接复制粘贴）", 13, True, INDIGO_DARK, first=True, space_after=5)
    for ln in prompt_lines:
        para(qtf, ln, 12, False, TEXT, space_after=2)
    obox = add_box(s, Inches(0.7), Inches(6.78), Inches(12.0), Inches(0.4))
    para(obox.text_frame, output, 14.5, True, AMBER, first=True)
    footer(s)
    return s


def two_col_slide(title, left_title, left_items, right_title, right_items, bottom=None,
                  subtitle=None):
    s = blank_slide()
    header(s, title, subtitle)
    lbox = add_box(s, Inches(0.7), Inches(1.62), Inches(5.9), Inches(4.6))
    para(lbox.text_frame, left_title, 16.5, True, INDIGO_DARK, first=True, space_after=8)
    for it in left_items:
        para(lbox.text_frame, "▪ " + it, 14, False, TEXT, space_after=6)
    rbox = add_box(s, Inches(6.85), Inches(1.62), Inches(5.9), Inches(4.6))
    para(rbox.text_frame, right_title, 16.5, True, INDIGO_DARK, first=True, space_after=8)
    for it in right_items:
        para(rbox.text_frame, "▪ " + it, 14, False, TEXT, space_after=6)
    if bottom:
        bb = add_rect(s, Inches(0.7), Inches(6.45), Inches(12.05), Inches(0.55),
                      AMBER_LIGHT, AMBER)
        rect_text(bb, bottom, 14, True, RGBColor(0x92, 0x40, 0x0E), PP_ALIGN.LEFT)
    footer(s)
    return s


# ============ P1 封面 ============
s = deep_slide()
add_box(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(0.5)).text_frame
kb = add_box(s, Inches(0.9), Inches(1.05), Inches(11.5), Inches(0.5))
para(kb.text_frame, "LUMITON·OMAX 用户培训", 18, True, RGBColor(0xA5, 0xB4, 0xFC), first=True)
tb = add_box(s, Inches(0.9), Inches(2.15), Inches(11.6), Inches(1.9))
para(tb.text_frame, "让每一轮研发", 44, True, WHITE, first=True, space_after=4)
para(tb.text_frame, "都可追踪、可复盘、可预测", 44, True, WHITE)
sub = add_box(s, Inches(0.9), Inches(4.15), Inches(11.6), Inches(0.6))
para(sub.text_frame, "—— 证据驱动的研发闭环工作法：用 Lumiton·Omax 建立可复制的研发操作系统",
     20, False, RGBColor(0xC7, 0xD2, 0xFE), first=True)
add_rect(s, Inches(0.92), Inches(4.95), Inches(1.4), Pt(3), AMBER, shape=MSO_SHAPE.RECTANGLE)
info = add_box(s, Inches(0.9), Inches(5.35), Inches(11.5), Inches(0.8))
para(info.text_frame, "对象：主管领导 + 系统使用人员（约 12 人）", 14, False,
     RGBColor(0xA5, 0xB4, 0xFC), first=True, space_after=2)
para(info.text_frame, "形式：PPT 宣讲为主线 · 试用系统现场实操 · 2026 年 8 月", 14, False,
     RGBColor(0xA5, 0xB4, 0xFC))

# ============ P2 议程 ============
s = blank_slide()
header(s, "今天的议程", "半天 3 小时 · 不是功能培训，是一起跑通一条研发链")
agenda = [
    ("一", "开场：我们从试用中观察到什么", "30 分钟", "全员"),
    ("二", "一条主线：证据驱动研发七阶段（含 4 次现场实操）", "100 分钟", "使用人员为主"),
    ("三", "把闭环装进系统：试点笔记本、命名规则、数据模板", "20 分钟", "使用人员"),
    ("四", "数据安全：出网脱敏对使用的影响", "10 分钟", "全员"),
    ("五", "稳定保障与路线图：从 7 月反馈到 8 月修复", "10 分钟", "全员"),
    ("六", "从今天开始的 30 天：成功指标与每个人的下一步", "10 分钟", "全员"),
]
y = 1.7
for tag, name, dur, who in agenda:
    add_rect(s, Inches(0.7), Inches(y), Inches(0.55), Inches(0.72), INDIGO)
    rect_text(s.shapes[-1], tag, 16, True, WHITE)
    nb = add_box(s, Inches(1.45), Inches(y + 0.06), Inches(8.6), Inches(0.62))
    para(nb.text_frame, name, 16, True, TEXT, first=True)
    db = add_box(s, Inches(10.2), Inches(y + 0.08), Inches(2.4), Inches(0.5))
    para(db.text_frame, dur + "  ·  " + who, 12.5, False, SUB, first=True)
    y += 0.82
footer(s)

# ============ P3 章节1 ============
section_slide("01", "开场：我们从试用中观察到什么",
              "从试用期真实数据出发，对齐一个判断：问题不在 AI，在研发方法。")

# ============ P4 观察一 ============
s = blank_slide()
header(s, "观察一：资料已经很多，但以「文档堆叠」为主")
items = [
    "产品说明书、研发报告、评价实验报告、合成实验记录、原始数据表、结题报告……",
    "10+ 产品线并存：降失水剂、缓凝剂、减阻剂、冲洗剂、隔离剂、堵漏剂、膨胀剂、填充剂、加重剂、防沉降剂",
    "数据形态混合：PDF / DOC / XLS / ZIP / 外部文献链接并存",
    "同一主题连续迭代：HX-16L / HX-16S 降失水剂、ATMP 合成条件优化……",
    "更偏「应用结果」：缺少统一的原料批次、分子结构、工艺参数、测试条件、失败原因字段",
]
lb = add_box(s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(4.4))
para(lb.text_frame, "▪ " + items[0], 15.5, False, TEXT, first=True, space_after=10)
for it in items[1:]:
    para(lb.text_frame, "▪ " + it, 15.5, False, TEXT, space_after=10)
bb = add_rect(s, Inches(0.7), Inches(6.45), Inches(12.05), Inches(0.55), AMBER_LIGHT, AMBER)
rect_text(bb, "问题不是资料不够 —— 是资料不可检索、不可复用。", 14, True,
          RGBColor(0x92, 0x40, 0x0E), PP_ALIGN.LEFT)
footer(s)

# ============ P5 观察二 ============
s = blank_slide()
header(s, "观察二：典型的提问方式是「直接要答案」")
items = [
    "把大量实验报告、产品说明书、测试记录导入系统后",
    "直接问：「下一步怎么做？」「给出建议」",
    "本质：把研发决策外包给 AI",
    "AI 的建议没有经过问题定义、证据整理、假设验证",
    "结果：建议不可追踪、不可复盘、更不可预测",
]
lb = add_box(s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(4.4))
para(lb.text_frame, "▪ " + items[0], 15.5, False, TEXT, first=True, space_after=10)
for it in items[1:]:
    para(lb.text_frame, "▪ " + it, 15.5, False, TEXT, space_after=10)
bb = add_rect(s, Inches(0.7), Inches(6.45), Inches(12.05), Inches(0.55), AMBER_LIGHT, AMBER)
rect_text(bb, "这不是提问技巧问题 —— 是研发活动缺一套可被 AI 承接的结构化方法。", 14,
          True, RGBColor(0x92, 0x40, 0x0E), PP_ALIGN.LEFT)
footer(s)

# ============ P6 观察三 ============
s = blank_slide()
header(s, "观察三：三个层级，其实在问三件不同的事")
cols = [
    ("管理层", ["降本 10%-15%", "现场技术问题 -30%", "交付率 85%+", "决策支持而非一次性报告"],
     "要「管理视图」"),
    ("中层", ["补上分子结构/理化性质验证", "解释水泥原料波动的影响", "总结配方-条件-性能规律"],
     "要「机理与预测」"),
    ("执行层", ["同一代号多配方、现场对不上", "高温高压测试反复试错", "失败经验不可复用"],
     "要「结构化资产」"),
]
x = 0.7
for title, its, want in cols:
    hd = add_rect(s, Inches(x), Inches(1.7), Inches(3.9), Inches(0.62), INDIGO)
    rect_text(hd, title, 17, True, WHITE)
    ib = add_box(s, Inches(x + 0.05), Inches(2.5), Inches(3.85), Inches(3.1))
    para(ib.text_frame, "▪ " + its[0], 13.5, False, TEXT, first=True, space_after=8)
    for it in its[1:]:
        para(ib.text_frame, "▪ " + it, 13.5, False, TEXT, space_after=8)
    wb = add_rect(s, Inches(x), Inches(5.75), Inches(3.9), Inches(0.55), AMBER_LIGHT, AMBER)
    rect_text(wb, want, 14, True, RGBColor(0x92, 0x40, 0x0E))
    x += 4.15
bb = add_rect(s, Inches(0.7), Inches(6.45), Inches(12.05), Inches(0.55), INDIGO_LIGHT, INDIGO)
rect_text(bb, "三层诉求不同，但都指向同一个缺口 —— 一条能被 AI 承接的研发链。", 14, True,
          INDIGO_DARK, PP_ALIGN.LEFT)
footer(s)

# ============ P7 核心判断 ============
s = blank_slide()
header(s, "我们的核心判断")
qb = add_rect(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(2.6), INDIGO_LIGHT, INDIGO)
qtf = qb.text_frame
qtf.margin_left = Pt(20)
qtf.margin_top = Pt(24)
para(qtf, "当前试用面临的主要问题，不是「不会问 AI」——", 22, True, INDIGO_DARK,
     first=True, space_after=8)
para(qtf, "而是研发活动本身，缺少一套可被 AI 承接的结构化方法。", 22, True, INDIGO_DARK)
vb = add_rect(s, Inches(0.9), Inches(4.85), Inches(11.5), Inches(1.55), AMBER_LIGHT, AMBER)
vtf = vb.text_frame
vtf.margin_left = Pt(20)
vtf.margin_top = Pt(14)
para(vtf, "Lumina 在欧美克的价值主张：", 15, True, RGBColor(0x92, 0x40, 0x0E), first=True,
     space_after=6)
para(vtf, "从「实验报告问答工具」→ 把历史经验、实验数据、现场反馈、产品机理假设", 16.5,
     True, TEXT)
para(vtf, "组织成可追踪、可复盘、可预测的研发决策系统", 16.5, True, TEXT)
footer(s)

# ============ P8 章节2 ============
section_slide("02", "一条主线：证据驱动研发七阶段",
              "每个阶段 = 一个研发动作 + 一组标准提问 + 一份可沉淀的产出。")

# ============ P9 七阶段总览 ============
s = blank_slide()
header(s, "七阶段研发闭环", "试点：油井水泥用降失水剂 —— 一线最高优先级，回应规律总结与预测诉求")
stages = [
    ("阶段 0", "问题定义", "立题与边界"),
    ("阶段 1", "证据地图", "分层与索引"),
    ("阶段 2", "产品-配方-场景映射", "代号-版本-场景"),
    ("阶段 3", "机理假设", "证据与反证"),
    ("阶段 4", "实验矩阵", "≤8 组最小矩阵"),
    ("阶段 5", "实验复盘", "假设 vs 结果"),
    ("阶段 6", "Gate 与管理摘要", "决策与管理视图"),
]
x = 0.7
w = 1.66
gap = 0.08
for i, (tag, name, sub) in enumerate(stages):
    fill = AMBER if i == 6 else INDIGO
    blk = add_rect(s, Inches(x), Inches(1.85), Inches(w), Inches(1.9), fill)
    tf = blk.text_frame
    tf.margin_top = Pt(10)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = tag
    set_run(r, 12.5, True, WHITE)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)
    r2 = p2.add_run()
    r2.text = name
    set_run(r2, 15.5, True, WHITE)
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(5)
    r3 = p3.add_run()
    r3.text = sub
    set_run(r3, 11.5, False, RGBColor(0xE0, 0xE7, 0xFF))
    if i < 6:
        ab = add_box(s, Inches(x + w + 0.008), Inches(2.62), Inches(gap), Inches(0.35))
        para(ab.text_frame, "›", 16, True, SUB, PP_ALIGN.CENTER, first=True)
    x += w + gap
mb = add_rect(s, Inches(0.7), Inches(4.1), Inches(12.05), Inches(0.62), INDIGO_LIGHT, INDIGO)
rect_text(mb, "每一阶段 = 一个研发动作 + 一组标准提问 + 一份可沉淀的产出（保存为笔记）",
          14.5, True, INDIGO_DARK, PP_ALIGN.LEFT)
nb = add_rect(s, Inches(0.7), Inches(4.95), Inches(12.05), Inches(0.62), AMBER_LIGHT, AMBER)
rect_text(nb, "每个阶段的「标准提问」已内置在系统提示与帮助中心 —— 忘了怎么问，系统会引导你",
          14, True, RGBColor(0x92, 0x40, 0x0E), PP_ALIGN.LEFT)
footer(s)

# ============ P10 阶段0 ============
two_col_slide(
    "阶段 0：项目立题与边界定义",
    "研发人员必须先回答",
    ["本项目服务哪个产品线？",
     "目标应用区块、温度、压力、水泥厂家/批次、水质条件？",
     "当前失败或不满足指标的具体表现？",
     "成功判据是什么？",
     "成本、工艺、交付周期有什么约束？",
     "已尝试过哪些方案？"],
    "Lumina 的任务",
    ["根据上传资料生成「项目问题定义卡」",
     "识别资料中缺失的关键边界条件",
     "把宽泛问题改写成可实验验证的问题"],
    "目标：防止一上来就「许愿」—— 先定义问题，再谈方案",
    "下一张：现场实操"
)

# ============ P11 实操1 ============
practice_slide(
    1, "0 ｜ 问题定义卡",
    ["新建试点笔记本「降失水剂_HX-16L_试点」",
     "上传 2–3 份主题资料（研发报告 / 实验记录 Excel / 产品说明书）",
     "进入 Chat，粘贴右侧提问词"],
    ["请基于当前笔记本资料，帮我生成本项目的问题定义卡。必须包含：",
     "1. 产品对象　2. 应用工况　3. 当前不达标指标　4. 已尝试方案",
     "5. 成功判据　6. 关键约束　7. 还缺哪些数据才能进入实验设计",
     "请把「文档中有证据」和「你推断的内容」分开列出。"],
    "产出：问题定义卡 → 保存为笔记「HX-16L 问题定义卡 v1」"
)

# ============ P12 阶段1 ============
two_col_slide(
    "阶段 1：资料分层与证据地图",
    "建议 5 类证据笔记本",
    ["01 产品说明与技术指标",
     "02 历史研发报告",
     "03 实验记录与原始数据",
     "04 现场反馈与施工结果",
     "05 外部文献与竞品资料"],
    "Lumina 的任务",
    ["汇总每份资料的证据类型",
     "标记「可用于决策」与「仅作背景」",
     "建立「指标—实验—结论—文件来源」证据索引",
     "输出每份资料的证据强度：高 / 中 / 低"],
    "目标：把一堆报告变成可查询的证据体系",
    "下一张：现场实操"
)

# ============ P13 实操2a ============
practice_slide(
    2, "1 ｜ 证据地图",
    ["在来源栏演示筛选与「本次聊天引用范围」选择",
     "粘贴右侧提问词"],
    ["请把当前来源按「产品说明、实验记录、研发报告、现场反馈、外部文献、原始数据」分类。",
     "对每个来源判断：",
     "1. 它能支持哪类研发决策　2. 它包含哪些关键指标",
     "3. 它缺少哪些字段　4. 它的证据强度是高/中/低",
     "最后输出一张证据地图。"],
    "产出：证据地图 → 保存为笔记"
)

# ============ P14 阶段2 ============
two_col_slide(
    "阶段 2：产品-配方-场景映射",
    "最小映射字段",
    ["产品代号：HX-12L / HX-14L / HX-16L / HX-16S",
     "配方版本：小试 / 中试 / 生产样 / 现场样",
     "功能类型 + 原料批次与供应商",
     "水泥厂家、批次、类型 · 水质条件",
     "工况：温度、压力、密度、盐度",
     "关键性能指标 · 现场或实验结果"],
    "Lumina 的任务",
    ["从历史报告抽取产品代号与配方版本线索",
     "发现同名产品下的配方差异",
     "识别「换水泥后性能变化」的案例"],
    "目标：解决「同一代号多配方、现场对不上」的核心管理问题",
    "下一张：现场实操（重点演示来源筛选）"
)

# ============ P15 实操2b ============
practice_slide(
    3, "2 ｜ 产品-配方-场景映射",
    ["来源栏筛选本次要参考的来源（如「HX」）—— 资料属于项目 ≠ 本次回答引用",
     "粘贴右侧提问词"],
    ["请从当前资料中整理所有降失水剂相关的产品代号、配方版本、应用场景和性能结果。",
     "重点识别：",
     "1. 同一产品代号是否对应多个配方或样品阶段",
     "2. 是否出现同一配方在不同水泥/水质下表现差异",
     "3. 哪些记录缺少配方、批次或测试条件，导致无法复用",
     "请输出为表格。"],
    "产出：产品-配方-场景映射表 → 保存为笔记"
)

# ============ P16 阶段3 ============
two_col_slide(
    "阶段 3：机理假设生成",
    "每个假设必须包含",
    ["假设描述",
     "已有证据",
     "可能反证",
     "需要补测的理化指标",
     "推荐验证实验",
     "若成立，下一轮配方调整方向"],
    "Lumina 的任务",
    ["基于历史报告 + 外部文献生成可验证假设",
     "把假设拆成「支持证据 / 反证风险 / 验证实验」",
     "避免把相关性直接说成因果",
     "回答区分「已有证据 / 合理推断 / 需要补充验证」"],
    "目标：补上分子结构验证 —— 不问「下一步怎么做」，先问「哪些机理可能解释差异」",
    "下一张：现场实操"
)

# ============ P17 实操3a ============
practice_slide(
    4, "3 ｜ 机理假设",
    ["基于映射表与证据地图继续对话",
     "粘贴右侧提问词"],
    ["请基于当前降失水剂资料，提出 3-5 个可能解释性能差异的机理假设。",
     "每个假设必须包含：",
     "1. 假设描述　2. 已有证据　3. 可能反证",
     "4. 需要补测的理化指标　5. 推荐验证实验　6. 若假设成立，下一轮配方调整方向",
     "不要直接给最终配方。"],
    "产出：机理假设清单（区分证据 / 推断 / 待验证）"
)

# ============ P18 阶段4 ============
two_col_slide(
    "阶段 4：实验矩阵设计",
    "每轮实验必须明确",
    ["固定变量 · 待考察变量",
     "对照组",
     "指标优先级",
     "失败判据",
     "是否进入中试的 Gate",
     "实验数量：不超过 8 组"],
    "Lumina 的任务",
    ["根据假设生成最小实验矩阵",
     "减少无效组合",
     "给出「为什么这样设计」的解释"],
    "目标：把 AI 建议变成可执行实验，而不是灵感列表",
    "下一张：现场实操"
)

# ============ P19 实操3b ============
practice_slide(
    5, "4 ｜ 实验矩阵",
    ["选定一个假设继续提问",
     "粘贴右侧提问词"],
    ["请围绕假设 A 设计一轮最小实验矩阵。要求：",
     "1. 控制实验数量，不超过 8 组　2. 明确每组改变的变量　3. 明确对照组",
     "4. 明确测试指标和通过标准　5. 明确失败后如何判断原因",
     "6. 说明哪些结果支持或推翻该假设。"],
    "产出：实验矩阵表（≤8 组、含失败判据与 Gate 条件）"
)

# ============ P20 阶段5 ============
two_col_slide(
    "阶段 5：实验后复盘",
    "复盘必须回答",
    ["原假设是什么",
     "哪些数据支持假设",
     "哪些数据与假设冲突",
     "失败最可能来自：配方 / 工艺 / 原料 / 水泥批次 / 测试条件",
     "下一轮只验证哪 1-2 个问题",
     "这次实验应沉淀为什么规则"],
    "Lumina 的任务",
    ["对比实验前假设与实验后结果",
     "判断假设被支持、部分支持还是推翻",
     "生成下一轮实验问题",
     "输出可保存为笔记的结构化复盘"],
    "目标：每次实验都沉淀，而不是只留下一个报告文件",
    "下一张：现场实操"
)

# ============ P21 实操4a ============
practice_slide(
    6, "5 ｜ 实验复盘",
    ["上传本轮实验报告作为新来源",
     "粘贴右侧提问词",
     "点击「保存为笔记」—— 每次实验必须沉淀，失败案例进入知识库"],
    ["请根据本轮实验报告，做一次研发复盘。必须回答：",
     "1. 原假设是什么　2. 哪些数据支持假设　3. 哪些数据与假设冲突",
     "4. 失败或异常最可能来自配方、工艺、原料、水泥批次还是测试条件",
     "5. 下一轮只应该验证哪 1-2 个问题",
     "6. 这次实验应该沉淀到知识库的规则是什么。"],
    "产出：结构化复盘笔记 → 失败案例复用从这里开始"
)

# ============ P22 阶段6 ============
two_col_slide(
    "阶段 6：项目 Gate 与管理层摘要（主管领导重点）",
    "4 个 Gate",
    ["Gate 1：问题定义完成，关键资料齐备",
     "Gate 2：机理假设明确，实验矩阵通过评审",
     "Gate 3：小试结果达标，成本与工艺风险可控",
     "Gate 4：中试/现场验证通过，进入产品化或客户应用"],
    "Lumina 的任务",
    ["生成每周项目状态摘要",
     "标记当前 Gate、风险和阻塞点",
     "输出继续 / 暂停 / 转向 / 补数据建议",
     "每个项目输出四类视图：成本影响 · 现场风险 · 交付进度 · 复用知识"],
    "目标：把技术过程转成能用于决策的管理视图；研发例会固定使用",
    "下一张：现场实操"
)

# ============ P23 实操4b ============
practice_slide(
    7, "6 ｜ 管理摘要",
    ["面向主管领导演示",
     "粘贴右侧提问词"],
    ["请把当前项目状态整理成给总经理和研发副总看的管理摘要。",
     "必须包含：",
     "1. 当前 Gate　2. 已达成的证据　3. 未解决的关键风险",
     "4. 对成本、现场风险、交付周期的影响　5. 本周需要管理层协调的事项",
     "6. 是否建议继续推进、补实验、暂停或调整方向。"],
    "产出：周报级管理摘要 → 研发例会固定使用"
)

# ============ P24 阶段衔接 ============
s = blank_slide()
header(s, "忘了标准提问？系统会主动引导你走闭环")
c1 = add_rect(s, Inches(0.7), Inches(1.8), Inches(5.95), Inches(3.9), INDIGO_LIGHT, INDIGO)
tf = c1.text_frame
tf.margin_left = Pt(16)
tf.margin_top = Pt(16)
para(tf, "导览卡片", 18, True, INDIGO_DARK, first=True, space_after=8)
para(tf, "首次导入来源后自动生成：", 14, True, TEXT, space_after=6)
para(tf, "▪ 笔记本摘要", 14, False, TEXT, space_after=5)
para(tf, "▪ 关键要点", 14, False, TEXT, space_after=5)
para(tf, "▪ 3 条建议问题（点击直接发送）", 14, False, TEXT)
c2 = add_rect(s, Inches(6.8), Inches(1.8), Inches(5.95), Inches(3.9), INDIGO_LIGHT, INDIGO)
tf = c2.text_frame
tf.margin_left = Pt(16)
tf.margin_top = Pt(16)
para(tf, "回答后的下一步建议", 18, True, INDIGO_DARK, first=True, space_after=8)
para(tf, "每轮主回答结束后给出 3 条建议问题：", 14, True, TEXT, space_after=6)
para(tf, "▪ 结合本轮问题 + 主回答 + 笔记本上下文", 14, False, TEXT, space_after=5)
para(tf, "▪ 主回答完成即恢复输入框，建议异步生成不阻塞", 14, False, TEXT, space_after=5)
para(tf, "▪ 连续多轮回答，建议各自跟随当轮 AI 消息", 14, False, TEXT)
bb = add_rect(s, Inches(0.7), Inches(6.0), Inches(12.05), Inches(0.62), AMBER_LIGHT, AMBER)
rect_text(bb, "两条内置引导 = 把「许愿式提问」拉回研发闭环", 14.5, True,
          RGBColor(0x92, 0x40, 0x0E), PP_ALIGN.LEFT)
footer(s)

# ============ P25 章节3 ============
section_slide("03", "把闭环装进系统",
              "试点笔记本体系、命名规则与最小数据模板 —— 30 分钟就能搭起来。")

# ============ P26 试点笔记本 ============
s = blank_slide()
header(s, "试点笔记本体系", "1 个聚合笔记本 + 6 个子笔记本")
top = add_rect(s, Inches(0.7), Inches(1.7), Inches(12.05), Inches(0.85), INDIGO_DEEP)
rect_text(top, "聚合笔记本：「降失水剂研发方法试点」—— 动态视图关联，不拷贝数据，源笔记本照常更新",
          15.5, True, WHITE)
subs = [
    "子笔记本 1\n降失水剂产品说明与指标",
    "子笔记本 2\nHX-12/HX-14/HX-16\n历史研发报告",
    "子笔记本 3\n降失水剂实验记录\n与原始数据",
    "子笔记本 4\n中试、生产样\n与现场反馈",
    "子笔记本 5\n水泥批次、水质\n与工况条件",
    "子笔记本 6\n外部文献与竞品资料",
]
x = 0.7
w = 1.95
gap = 0.07
for t in subs:
    blk = add_rect(s, Inches(x), Inches(2.85), Inches(w), Inches(1.75), INDIGO)
    tf = blk.text_frame
    tf.margin_top = Pt(8)
    for j, ln in enumerate(t.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = ln
        set_run(r, 12.5, j == 0, WHITE if j == 0 else RGBColor(0xE0, 0xE7, 0xFF))
    x += w + gap
nb = add_box(s, Inches(0.7), Inches(4.9), Inches(12.0), Inches(0.5))
para(nb.text_frame, "命名建议：产品线_产品代号_样品阶段_应用条件_日期_资料类型", 15, True,
     INDIGO_DARK, first=True)
mb = add_rect(s, Inches(0.7), Inches(5.5), Inches(12.05), Inches(0.95), AMBER_LIGHT, AMBER)
mtf = mb.text_frame
mtf.margin_left = Pt(14)
para(mtf, "聚合视图里怎么查：来源/笔记均可按子笔记本、上传人、标题筛选；每个子笔记本独立维护，聚合自动汇总。",
     13.5, False, RGBColor(0x92, 0x40, 0x0E), first=True)
footer(s)

# ============ P27 命名规则 ============
s = blank_slide()
header(s, "来源命名规则", "上传文件时保留原始文件名；标题按规则规范，便于筛选与映射")
examples = [
    "降失水剂_HX-16L_中试样_中低温防沉降_2026-05-09_评价报告",
    "降失水剂_HX-16L_小试样_合成记录_2025-09-12_实验记录",
    "降失水剂_HX-14L_优化评价_2025-02-08_实验报告",
]
y = 1.9
for ex in examples:
    eb = add_rect(s, Inches(0.7), Inches(y), Inches(12.05), Inches(0.62), INDIGO_LIGHT,
                  INDIGO)
    rect_text(eb, ex, 14.5, False, INDIGO_DARK, PP_ALIGN.LEFT)
    y += 0.78
rules = [
    "同一产品代号必须对应到具体配方版本（小试/中试/生产/现场），否则无法复用",
    "上传重复文件会被识别（忽略大小写与首尾空格），只上传非重复项",
    "下载文件自动还原原始文件名；同名版本靠标题时间戳区分",
]
rb = add_box(s, Inches(0.7), Inches(4.35), Inches(12.0), Inches(1.8))
para(rb.text_frame, "▪ " + rules[0], 14.5, False, TEXT, first=True, space_after=8)
for r in rules[1:]:
    para(rb.text_frame, "▪ " + r, 14.5, False, TEXT, space_after=8)
footer(s)

# ============ P28 最小数据模板 ============
s = blank_slide()
header(s, "最小数据模板：实验记录", "上传报告之外，补一张结构化表；初期不追求完美数据库")
left = [
    "项目名称 — 如中低温防沉降降失水剂研发",
    "产品代号 — 如 HX-16L",
    "配方版本 — 小试样 / 中试样 / 生产样 / 现场样",
    "功能类型 — 降失水 / 缓凝 / 减阻 / 复配",
    "原料与批次 — 关键单体、聚合物、供应商、批号",
    "合成工艺参数 — 温度、时间、pH、引发剂、加料方式",
    "理化指标 — 分子量、分布、固含、黏度、zeta 电位、热稳定性",
    "水泥条件 — 厂家、批次、G 级水泥类型、储存状态",
]
right = [
    "水质条件 — 淡水 / 盐水 / 离子浓度",
    "工况条件 — 温度、压力、密度、盐度",
    "测试指标 — 降失水、流变、稠化、强度、稳定性",
    "结果判定 — 达标 / 部分达标 / 不达标",
    "异常现象 — 包芯、倒挂、沉降、失稳、强度不足",
    "失败原因假设 — 配方 / 工艺 / 水泥 / 水质 / 测试 / 未知",
    "下一步动作 — 补测、调整配方、复测、中试、暂停",
]
lb = add_box(s, Inches(0.7), Inches(1.75), Inches(6.0), Inches(4.4))
para(lb.text_frame, "▪ " + left[0], 13.5, False, TEXT, first=True, space_after=7)
for it in left[1:]:
    para(lb.text_frame, "▪ " + it, 13.5, False, TEXT, space_after=7)
rb2 = add_box(s, Inches(6.95), Inches(1.75), Inches(5.9), Inches(4.4))
para(rb2.text_frame, "▪ " + right[0], 13.5, False, TEXT, first=True, space_after=7)
for it in right[1:]:
    para(rb2.text_frame, "▪ " + it, 13.5, False, TEXT, space_after=7)
bb = add_rect(s, Inches(0.7), Inches(6.45), Inches(12.05), Inches(0.55), AMBER_LIGHT, AMBER)
rect_text(bb, "字段越完整，阶段 2 的映射与阶段 3 的假设就越有依据 —— 这是「预测」的前提", 14,
          True, RGBColor(0x92, 0x40, 0x0E), PP_ALIGN.LEFT)
footer(s)

# ============ P29 现场反馈模板 ============
s = blank_slide()
header(s, "最小数据模板：现场反馈", "现场失败代价最高 —— 每条现场反馈都要沉淀")
left = [
    "客户 / 区块 — 现场应用对象",
    "井号 / 工况 — 井深、温度、压力、施工条件",
    "使用产品 — 产品代号、配方版本、生产批次",
    "水泥与水质 — 厂家、批次、水质条件",
]
right = [
    "施工表现 — 泵送、稠化、返高、候凝、异常",
    "固井质量 — 合格 / 问题描述",
    "客户反馈 — 原话或结构化摘要",
    "研发判断 — 与配方 / 工况 / 水泥波动的关系",
    "后续动作 — 复测、调整、跟踪、关闭",
]
lb = add_box(s, Inches(0.7), Inches(1.9), Inches(6.0), Inches(3.6))
para(lb.text_frame, "▪ " + left[0], 14.5, False, TEXT, first=True, space_after=10)
for it in left[1:]:
    para(lb.text_frame, "▪ " + it, 14.5, False, TEXT, space_after=10)
rb2 = add_box(s, Inches(6.95), Inches(1.9), Inches(5.9), Inches(3.6))
para(rb2.text_frame, "▪ " + right[0], 14.5, False, TEXT, first=True, space_after=10)
for it in right[1:]:
    para(rb2.text_frame, "▪ " + it, 14.5, False, TEXT, space_after=10)
bb = add_rect(s, Inches(0.7), Inches(6.0), Inches(12.05), Inches(0.62), INDIGO_LIGHT, INDIGO)
rect_text(bb, "现场反馈 → 回到阶段 0/2 重新核对问题定义与产品-配方-场景映射", 14.5, True,
          INDIGO_DARK, PP_ALIGN.LEFT)
footer(s)

# ============ P30 脱敏 ============
s = blank_slide()
header(s, "数据安全：出网脱敏，使用无感", "库内数据不动 · 只在发往外部 AI 时临时替换 · 返回自动还原")
top = add_rect(s, Inches(0.7), Inches(1.7), Inches(12.05), Inches(0.95), INDIGO_DEEP)
tft = top.text_frame
tft.margin_left = Pt(14)
para(tft, "透明网关原理：库内原文（笔记本/来源/笔记/向量/图谱）始终保持原始形态，", 14.5,
     False, WHITE, first=True, space_after=2)
para(tft, "仅在内容发往外部 AI 模型瞬间做一次性替换，模型返回后自动还原 —— 您看到的界面与原来一致。",
     14.5, False, WHITE)
notes = [
    ("注意 1：电话显示 888888", "按规则固化不可还原，这是设计行为；需要真实电话请查阅原始资料。"),
    ("注意 2：联网搜索用科学主题词", "身份词（人名/井号）以代号形态出网，公网无法识别；请用「高温下降失水剂性能」这类主题词。"),
    ("注意 3：新人名需录词典", "井号/电话/产品代号自动识别；人员姓名需管理员在设置页录入词典后才会被替换。"),
]
x = 0.7
for title, body in notes:
    blk = add_rect(s, Inches(x), Inches(3.0), Inches(3.95), Inches(2.4), AMBER_LIGHT, AMBER)
    tf = blk.text_frame
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(12)
    para(tf, title, 14.5, True, RGBColor(0x92, 0x40, 0x0E), first=True, space_after=8)
    para(tf, body, 12.5, False, TEXT)
    x += 4.05
bb = add_rect(s, Inches(0.7), Inches(5.8), Inches(12.05), Inches(0.62), INDIGO_LIGHT, INDIGO)
rect_text(bb, "库内存储、搜索、引用全部使用原文 —— 用真实人名/井号提问可正常命中，无需改变习惯", 14,
          True, INDIGO_DARK, PP_ALIGN.LEFT)
footer(s)

# ============ P31 稳定保障 ============
s = blank_slide()
header(s, "稳定保障：从 7 月反馈到 8 月修复")
lb = add_box(s, Inches(0.7), Inches(1.6), Inches(6.0), Inches(4.8))
para(lb.text_frame, "已修复（8 月）", 16.5, True, INDIGO_DARK, first=True, space_after=8)
para(lb.text_frame, "▪ 全局提问偶发「供应商返回错误」— 实为内部并行检索合并问题，已修复，与供应商无关", 13.5,
     False, TEXT, space_after=7)
para(lb.text_frame, "▪ 老版 .doc 上传后嵌入/图谱卡住 — 自动引擎路径已补转换，两份失败文档已恢复全流程", 13.5,
     False, TEXT, space_after=7)
para(lb.text_frame, "▪ 设置页密钥与白名单被误清空 — 改为掩码显示，保存其他设置不再误清", 13.5,
     False, TEXT, space_after=7)
para(lb.text_frame, "▪ 「询问与搜索 / 提问与搜索」名称统一", 13.5, False, TEXT, space_after=14)
para(lb.text_frame, "可观测（用户可见）", 16.5, True, INDIGO_DARK, space_after=8)
para(lb.text_frame, "▪ 阶段状态提示：获取上下文 / 联网搜索 / 生成回答 / 生成建议", 13.5, False,
     TEXT, space_after=7)
para(lb.text_frame, "▪ 等待秒数心跳：「正在等待模型响应（N 秒）」；超时以气泡提示操作指引", 13.5,
     False, TEXT, space_after=7)
para(lb.text_frame, "▪ 帮助中心 /help 全中文操作手册，随功能同步更新", 13.5, False, TEXT)
rb2 = add_box(s, Inches(6.95), Inches(1.6), Inches(5.9), Inches(4.8))
para(rb2.text_frame, "路线图（规划中，如实说明）", 16.5, True, RGBColor(0x92, 0x40, 0x0E),
     first=True, space_after=8)
para(rb2.text_frame, "▪ 模型故障自动切换机制", 13.5, False, TEXT, space_after=7)
para(rb2.text_frame, "▪ 实验数据索引（Excel 逐条结构化）", 13.5, False, TEXT, space_after=7)
para(rb2.text_frame, "▪ 实验时间线（按产品追溯）", 13.5, False, TEXT, space_after=7)
para(rb2.text_frame, "▪ 产品图谱与版本迭代图", 13.5, False, TEXT, space_after=7)
para(rb2.text_frame, "▪ 数据备份与恢复流程文档化", 13.5, False, TEXT, space_after=14)
para(rb2.text_frame, "现场保障", 16.5, True, INDIGO_DARK, space_after=8)
para(rb2.text_frame, "▪ 培训期间主备模型双保险", 13.5, False, TEXT, space_after=7)
para(rb2.text_frame, "▪ 管理员每周日志巡检与模型健康检查", 13.5, False, TEXT)
footer(s)

# ============ P32 成功指标 ============
s = blank_slide()
header(s, "成功指标：试点看得见", "领导看指标，团队看闭环")
cols = [
    ("试点 1 个月", [
        "试点聚合笔记本建成",
        "≥30 份历史资料形成证据地图",
        "产品代号-配方版本映射表",
        "3 次结构化问答沉淀为笔记",
        "1 份下一轮实验矩阵"], INDIGO),
    ("试点 3 个月", [
        "实验记录最小数据库",
        "3-5 条可复用研发规则",
        "减少重复试错实验组合",
        "≥1 个可进中试/现场验证方案",
        "可迁移到缓凝剂/减阻剂的流程模板"], INDIGO),
    ("管理层指标", [
        "关键项目按时交付率",
        "现场技术问题数量",
        "单产品研发成本",
        "失败案例复用率",
        "实验数据字段完整率",
        "问题→实验矩阵的周期"], AMBER),
]
x = 0.7
for title, its, color in cols:
    hd = add_rect(s, Inches(x), Inches(1.6), Inches(3.95), Inches(0.6), color)
    rect_text(hd, title, 15.5, True, WHITE)
    ib = add_box(s, Inches(x + 0.05), Inches(2.35), Inches(3.9), Inches(4.1))
    para(ib.text_frame, "▪ " + its[0], 13, False, TEXT, first=True, space_after=9)
    for it in its[1:]:
        para(ib.text_frame, "▪ " + it, 13, False, TEXT, space_after=9)
    x += 4.05
footer(s)

# ============ P33 下一步 ============
s = blank_slide()
header(s, "每个人的下一步")
cols = [
    ("主管领导", [
        "确定试点唯一口径：油井水泥用降失水剂",
        "明确 4 个 Gate 的审批节奏",
        "研发例会固定使用 Lumina 管理摘要"], INDIGO),
    ("研发人员", [
        "按命名规则整理资料，去重后上传",
        "建试点笔记本，用标准提问跑一轮闭环",
        "拒绝「许愿式提问」，先定义问题再问 AI",
        "每次实验后必做复盘并保存为笔记"], INDIGO),
    ("管理员", [
        "账号、模型与密钥配置保持稳定",
        "维护脱敏人名词典",
        "每周日志巡检 + 模型健康检查"], AMBER),
]
x = 0.7
for title, its, color in cols:
    hd = add_rect(s, Inches(x), Inches(1.6), Inches(3.95), Inches(0.6), color)
    rect_text(hd, title, 15.5, True, WHITE)
    ib = add_box(s, Inches(x + 0.05), Inches(2.35), Inches(3.9), Inches(4.1))
    para(ib.text_frame, "▪ " + its[0], 13.5, False, TEXT, first=True, space_after=10)
    for it in its[1:]:
        para(ib.text_frame, "▪ " + it, 13.5, False, TEXT, space_after=10)
    x += 4.05
bb = add_rect(s, Inches(0.7), Inches(6.3), Inches(12.05), Inches(0.62), INDIGO_LIGHT, INDIGO)
rect_text(bb, "本周内完成：试点笔记本建好 + 第一张问题定义卡", 14.5, True, INDIGO_DARK,
          PP_ALIGN.LEFT)
footer(s)

# ============ P34 结尾 ============
s = deep_slide()
tb = add_box(s, Inches(0.9), Inches(2.3), Inches(11.6), Inches(1.8))
para(tb.text_frame, "Lumina 的价值，不在回答问题，", 36, True, WHITE, first=True,
     space_after=6)
para(tb.text_frame, "而在帮助欧美克把研发过程变成可复制的操作系统。", 36, True, WHITE)
add_rect(s, Inches(0.92), Inches(4.3), Inches(1.4), Pt(3), AMBER, shape=MSO_SHAPE.RECTANGLE)
sb = add_box(s, Inches(0.9), Inches(4.7), Inches(11.6), Inches(1.0))
para(sb.text_frame, "从今天开始 —— 用一条闭环，让每一轮研发都可追踪、可复盘、可预测", 20,
     False, RGBColor(0xC7, 0xD2, 0xFE), first=True)
qb = add_box(s, Inches(0.9), Inches(6.0), Inches(11.6), Inches(0.6))
para(qb.text_frame, "现场答疑 · 提问词见各实操页 · 会后手册：系统内 /help 帮助中心", 14, False,
     RGBColor(0xA5, 0xB4, 0xFC), first=True)

OUT = "/Users/omax/YinShiApp/lumina-omax/docs/8-CUSTOMIZATION/2026-08-26-欧美克用户培训-研发操作系统.pptx"
prs.save(OUT)
print(f"saved: {OUT} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
